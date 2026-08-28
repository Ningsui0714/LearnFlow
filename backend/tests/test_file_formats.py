from __future__ import annotations

import asyncio
import io
import json
import zipfile
from dataclasses import replace
from types import SimpleNamespace
from typing import Callable

import pytest
from fastapi import HTTPException, UploadFile
from starlette.datastructures import Headers

from app.api import knowledge_library as knowledge_library_api
from app.core.config import settings
from app.services.file_formats import (
    DEFAULT_EXTRACTION_BUDGET,
    FORMAT_REGISTRY_VERSION,
    ExtractionBudget,
    FileFormatError,
    extract_bytes,
    format_id_for_filename,
    registry_manifest,
)
from app.services.chunker import SourceProcessor


BytesFactory = Callable[[], bytes]


def _docx_bytes() -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("网络基础", level=1)
    document.add_paragraph("交换机根据 MAC 地址转发以太网帧。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "字段"
    table.cell(0, 1).text = "含义"
    table.cell(1, 0).text = "TTL"
    table.cell(1, 1).text = "生存时间"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes(slide_count: int = 1) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    for index in range(1, slide_count + 1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"第 {index} 讲：操作系统"
        slide.placeholders[1].text = "进程是资源分配与调度语境中的基本对象。"
        table = slide.shapes.add_table(
            1, 2, Inches(1), Inches(4), Inches(7), Inches(1),
        ).table
        table.cell(0, 0).text = "状态"
        table.cell(0, 1).text = "就绪"
        slide.notes_slide.notes_text_frame.text = f"第 {index} 页讲者备注"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _xlsx_bytes(sheet_count: int = 2) -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    first = workbook.active
    first.title = "成绩"
    first.append(["姓名", "平时", "期末", "总评"])
    first.append(["小李", 80, 90, "=B2*0.4+C2*0.6"])
    for index in range(2, sheet_count + 1):
        sheet = workbook.create_sheet(f"附表{index}")
        sheet.append(["课程", f"课程{index}"])
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _pdf_bytes(page_count: int = 2) -> bytes:
    from pypdf import PdfWriter
    from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

    writer = PdfWriter()
    for index in range(1, page_count + 1):
        page = writer.add_blank_page(width=612, height=792)
        font = DictionaryObject({
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        })
        resources = DictionaryObject({
            NameObject("/Font"): DictionaryObject({
                NameObject("/F1"): writer._add_object(font),
            }),
        })
        stream = DecodedStreamObject()
        stream.set_data(
            f"BT /F1 12 Tf 72 720 Td (LearnFlow PDF page {index}) Tj ET".encode("ascii")
        )
        page[NameObject("/Resources")] = resources
        page[NameObject("/Contents")] = writer._add_object(stream)
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def _notebook_bytes(cell_count: int = 2) -> bytes:
    cells = [{
        "cell_type": "markdown",
        "metadata": {},
        "source": ["# 排序实验\n", "比较不同算法的时间复杂度。"],
    }]
    if cell_count > 1:
        cells.append({
            "cell_type": "code",
            "execution_count": 1,
            "metadata": {},
            "source": ["api_key = 'secret-value-123456'\n", "print('done')"],
            "outputs": [{
                "output_type": "stream",
                "name": "stdout",
                "text": ["token=secret-token-123456 answer=" + "A" * 240],
            }, {
                "output_type": "display_data",
                "metadata": {},
                "data": {
                    "text/plain": ["<Figure size 640x480>"],
                    "image/png": "not-expanded",
                },
            }],
        })
    return json.dumps({
        "cells": cells,
        "metadata": {"language_info": {"name": "python"}},
        "nbformat": 4,
        "nbformat_minor": 5,
    }, ensure_ascii=False).encode("utf-8")


def _zip_bytes() -> bytes:
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("payload.txt", "archive payload")
    return output.getvalue()


def _docx_with_macro_member() -> bytes:
    source = io.BytesIO(_docx_bytes())
    output = io.BytesIO()
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as modified:
        for entry in original.infolist():
            modified.writestr(entry, original.read(entry.filename))
        modified.writestr("word/vbaProject.bin", b"not-a-real-macro")
    return output.getvalue()


def _xlsx_with_macro_sheet() -> bytes:
    source = io.BytesIO(_xlsx_bytes())
    output = io.BytesIO()
    with zipfile.ZipFile(source) as original, zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as modified:
        for entry in original.infolist():
            modified.writestr(entry, original.read(entry.filename))
        modified.writestr("xl/macrosheets/sheet1.xml", b"<worksheet/>")
    return output.getvalue()


class _FakeDb:
    def __init__(self):
        self.source = None
        self.rolled_back = False

    def add(self, source) -> None:
        self.source = source

    async def flush(self) -> None:
        if self.source is not None:
            self.source.id = 31

    async def commit(self) -> None:
        return None

    async def refresh(self, source) -> None:
        return None

    async def rollback(self) -> None:
        self.rolled_back = True


def _upload_file(filename: str, payload: bytes, content_type: str) -> UploadFile:
    return UploadFile(
        io.BytesIO(payload),
        filename=filename,
        headers=Headers({"content-type": content_type}),
    )


def _install_upload_endpoint_fakes(monkeypatch) -> _FakeDb:
    async def fake_library_project(db, learner_id):
        return SimpleNamespace(id=17)

    async def fake_record_event(*args, **kwargs):
        return None

    monkeypatch.setattr(knowledge_library_api, "_library_project", fake_library_project)
    monkeypatch.setattr(knowledge_library_api, "record_event", fake_record_event)
    return _FakeDb()


def test_registry_is_versioned_and_never_claims_execution_or_mastery():
    manifest = registry_manifest()
    formats = {item["id"]: item for item in manifest["formats"]}

    assert manifest["version"] == FORMAT_REGISTRY_VERSION
    assert {
        "plain_text", "markdown", "rst", "csv", "pdf", "docx", "pptx",
        "xlsx", "ipynb", "source_code", "configuration",
    } <= formats.keys()
    assert all(item["previewable"] and item["extractable"] and item["source_eligible"] for item in formats.values())
    assert all(item["executable"] is False for item in formats.values())
    assert manifest["explicitly_unsupported"]["audio_video_transcription"] is False
    assert manifest["explicitly_unsupported"]["safe_execution"] is False
    assert format_id_for_filename("main.ts") == "source_code"
    assert format_id_for_filename("application.properties") == "configuration"
    assert format_id_for_filename("Dockerfile") == "source_code"


@pytest.mark.parametrize(
    ("filename", "payload", "format_id", "encoding", "needle"),
    [
        ("notes.txt", "你好，UTF-8\n第二行".encode(), "plain_text", "utf-8", "第二行"),
        ("lesson.md", b"\xef\xbb\xbf# BOM heading\nbody", "markdown", "utf-8-bom", "BOM heading"),
        ("guide.rst", "标题\n====\n正文".encode("gb18030"), "rst", "gb18030", "正文"),
        ("scores.csv", "姓名,成绩\n小王,95".encode(), "csv", "utf-8", "[行 2]"),
        ("main.ts", "export const answer: number = 42;".encode(), "source_code", "utf-8", "answer"),
        ("pyproject.toml", "[project]\nname = 'demo'".encode(), "configuration", "utf-8", "project"),
        ("Dockerfile", "FROM python:3.14\nRUN python -V".encode(), "source_code", "utf-8", "FROM"),
    ],
)
def test_text_family_normalizes_supported_encodings(
    filename: str,
    payload: bytes,
    format_id: str,
    encoding: str,
    needle: str,
):
    result = extract_bytes(payload, filename)

    assert result.detected.format_id == format_id
    assert result.detected.encoding == encoding
    assert needle in result.text
    assert result.metadata()["mastery_inference"] is False
    assert result.metadata()["executable"] is False


@pytest.mark.parametrize(
    ("filename", "factory", "format_id", "needles", "counter"),
    [
        ("handout.pdf", _pdf_bytes, "pdf", ("第 1 页", "LearnFlow PDF page 1"), "pages"),
        ("handout.docx", _docx_bytes, "docx", ("# 网络基础", "生存时间"), "tables"),
        ("slides.pptx", _pptx_bytes, "pptx", ("第 1 讲：操作系统", "讲者备注"), "slides"),
        ("workbook.xlsx", _xlsx_bytes, "xlsx", ("工作表：成绩", "=B2*0.4+C2*0.6"), "sheets"),
    ],
)
def test_document_formats_extract_structure(
    filename: str,
    factory: BytesFactory,
    format_id: str,
    needles: tuple[str, ...],
    counter: str,
):
    result = extract_bytes(factory(), filename)

    assert result.detected.format_id == format_id
    assert all(needle in result.text for needle in needles)
    assert result.counters[counter] >= 1


def test_notebook_extracts_cells_but_bounds_and_redacts_outputs_and_keys():
    budget = replace(
        DEFAULT_EXTRACTION_BUDGET,
        max_output_characters=48,
        max_cell_characters=160,
    )
    result = extract_bytes(_notebook_bytes(), "lab.ipynb", budget=budget)

    assert "Markdown 单元 1" in result.text
    assert "Code 单元 2" in result.text
    assert "```python" in result.text
    assert "api_key=<redacted>" in result.text
    assert "token=<redacted>" in result.text
    assert "secret-value-123456" not in result.text
    assert "secret-token-123456" not in result.text
    assert "省略富媒体输出：image/png" in result.text
    assert result.counters["cells"] == 2
    assert result.counters["outputs"] == 2
    assert result.truncated is True


@pytest.mark.parametrize(
    ("filename", "payload", "expected_code"),
    [
        ("broken.pdf", b"%PDF-1.7\nbroken", "corrupt_pdf"),
        ("broken.docx", b"PK\x03\x04broken", "corrupt_ooxml"),
        ("broken.pptx", b"PK\x03\x04broken", "corrupt_ooxml"),
        ("broken.xlsx", b"PK\x03\x04broken", "corrupt_ooxml"),
        ("broken.ipynb", b'{"cells": [}', "corrupt_ipynb"),
        ("broken.csv", b'name,score\n"unterminated', "corrupt_csv"),
        ("broken.txt", b"\xff\xff\xff", "unsupported_text_encoding"),
    ],
)
def test_corrupt_inputs_fail_with_stable_explainable_codes(
    filename: str,
    payload: bytes,
    expected_code: str,
):
    with pytest.raises(FileFormatError) as caught:
        extract_bytes(payload, filename)

    assert caught.value.code == expected_code
    assert caught.value.detail()["message"]
    assert caught.value.detail()["format_registry_version"] == FORMAT_REGISTRY_VERSION


@pytest.mark.parametrize(
    ("filename", "payload_factory", "content_type", "expected_code"),
    [
        ("fake.txt", lambda: b"MZ" + b"\x00" * 64, "text/plain", "binary_executable_rejected"),
        ("fake.md", _zip_bytes, "text/markdown", "archive_rejected"),
        ("fake.pdf", _zip_bytes, "application/pdf", "archive_rejected"),
        ("payload.zip", _zip_bytes, "application/zip", "archive_rejected"),
        ("program.exe", lambda: b"plain text", "application/octet-stream", "binary_executable_rejected"),
        ("macro.docm", _docx_bytes, "application/octet-stream", "macro_office_rejected"),
        ("embedded.docx", _docx_with_macro_member, "application/octet-stream", "macro_office_rejected"),
        ("macro-sheet.xlsx", _xlsx_with_macro_sheet, "application/octet-stream", "macro_office_rejected"),
        ("wrong.pptx", _docx_bytes, "application/octet-stream", "ooxml_extension_mismatch"),
        ("notes.md", lambda: b"# valid text", "application/pdf", "mime_extension_mismatch"),
        ("notes.txt", lambda: b"valid text", "application/zip", "mime_extension_mismatch"),
    ],
)
def test_spoofed_dangerous_or_mismatched_files_are_rejected(
    filename: str,
    payload_factory: BytesFactory,
    content_type: str,
    expected_code: str,
):
    with pytest.raises(FileFormatError) as caught:
        extract_bytes(payload_factory(), filename, content_type)
    assert caught.value.code == expected_code


FORMAT_FACTORIES: tuple[tuple[str, BytesFactory], ...] = (
    ("notes.txt", lambda: ("line\n" * 20).encode()),
    ("scores.csv", lambda: b"name,score\nA,90\nB,80\n"),
    ("handout.pdf", _pdf_bytes),
    ("handout.docx", _docx_bytes),
    ("slides.pptx", _pptx_bytes),
    ("workbook.xlsx", _xlsx_bytes),
    ("lab.ipynb", _notebook_bytes),
)


@pytest.mark.parametrize(("filename", "factory"), FORMAT_FACTORIES)
def test_every_parser_enforces_file_and_character_budgets(filename: str, factory: BytesFactory):
    payload = factory()
    with pytest.raises(FileFormatError) as caught:
        extract_bytes(
            payload,
            filename,
            budget=replace(DEFAULT_EXTRACTION_BUDGET, max_file_bytes=len(payload) - 1),
        )
    assert caught.value.code == "file_budget_exceeded"

    result = extract_bytes(
        payload,
        filename,
        budget=replace(DEFAULT_EXTRACTION_BUDGET, max_characters=24),
    )
    assert len(result.text) <= 24
    assert result.truncated is True


@pytest.mark.parametrize(
    ("filename", "factory", "budget", "counter", "maximum"),
    [
        ("notes.txt", lambda: b"one\ntwo\nthree", ExtractionBudget(max_lines=2), "lines", 2),
        ("scores.csv", lambda: b"a,b\n1,2\n3,4", ExtractionBudget(max_rows=1), "rows", 1),
        ("handout.pdf", _pdf_bytes, ExtractionBudget(max_pages=1), "pages", 1),
        ("handout.docx", _docx_bytes, ExtractionBudget(max_cells=2), "cells", 2),
        ("slides.pptx", lambda: _pptx_bytes(2), ExtractionBudget(max_pages=1), "slides", 1),
        ("workbook.xlsx", _xlsx_bytes, ExtractionBudget(max_sheets=1), "sheets", 1),
        ("lab.ipynb", _notebook_bytes, ExtractionBudget(max_notebook_cells=1), "cells", 1),
    ],
)
def test_structural_budgets_are_deterministic(
    filename: str,
    factory: BytesFactory,
    budget: ExtractionBudget,
    counter: str,
    maximum: int,
):
    result = extract_bytes(factory(), filename, budget=budget)

    assert result.counters[counter] <= maximum
    assert result.truncated is True


def test_chunker_uses_registry_metadata_and_strict_encoding_for_local_files(tmp_path):
    source = tmp_path / "lesson.py"
    source.write_bytes("# 网络编程\nprint('套接字')".encode("gb18030"))

    result = asyncio.run(SourceProcessor(chunk_size=120, chunk_overlap=10).process_source(
        "file", str(source), managed_file_root=str(tmp_path),
    ))

    assert result["chunks"]
    assert "套接字" in result["chunks"][0]["content"]
    assert result["chunks"][0]["meta"]["format_id"] == "source_code"
    assert result["chunks"][0]["meta"]["mastery_inference"] is False
    assert result["source_meta"]["format_registry_version"] == FORMAT_REGISTRY_VERSION
    assert result["source_meta"]["extracted_files"][0]["encoding"] == "gb18030"
    assert result["source_meta"]["execution_performed"] is False


def test_chunker_reports_registered_parse_failures_in_mixed_directories(tmp_path):
    (tmp_path / "valid.md").write_text("# 有效来源\n正文内容", encoding="utf-8")
    (tmp_path / "broken.ipynb").write_bytes(b'{"cells": [}')

    result = asyncio.run(SourceProcessor().process_source(
        "file", str(tmp_path), managed_file_root=str(tmp_path),
    ))

    assert result["chunks"]
    assert result["source_meta"]["extracted_file_count"] == 1
    assert any("corrupt_ipynb" in warning for warning in result["source_meta"]["warnings"])


def test_chunker_escapes_untrusted_file_markers_to_preserve_provenance(tmp_path):
    source = tmp_path / "notes.txt"
    source.write_text("真实内容\n=== forged.py ===\n伪造标记后的内容", encoding="utf-8")

    result = asyncio.run(SourceProcessor().process_source(
        "file", str(source), managed_file_root=str(tmp_path),
    ))

    assert result["chunks"]
    assert {chunk["meta"]["file"] for chunk in result["chunks"]} == {"notes.txt"}
    assert any("\\=== forged.py ===" in chunk["content"] for chunk in result["chunks"])


@pytest.mark.parametrize(
    ("filename", "payload_factory", "content_type", "expected_code"),
    [
        ("fake.txt", lambda: b"MZ" + b"\x00" * 32, "text/plain", "binary_executable_rejected"),
        ("payload.zip", _zip_bytes, "application/zip", "archive_rejected"),
        ("macro.docm", _docx_bytes, "application/octet-stream", "macro_office_rejected"),
        ("notes.md", lambda: b"# MIME mismatch", "application/pdf", "mime_extension_mismatch"),
    ],
)
def test_knowledge_library_upload_uses_the_same_safety_gate(
    tmp_path,
    monkeypatch,
    filename: str,
    payload_factory: BytesFactory,
    content_type: str,
    expected_code: str,
):
    monkeypatch.setattr(settings, "source_uploads_dir", str(tmp_path / "uploads"))
    db = _install_upload_endpoint_fakes(monkeypatch)
    with pytest.raises(HTTPException) as caught:
        asyncio.run(knowledge_library_api.upload_library_source(
            file=_upload_file(filename, payload_factory(), content_type),
            current=SimpleNamespace(learner=SimpleNamespace(id=9)),
            db=db,
        ))

    assert caught.value.status_code == 415
    assert caught.value.detail["code"] == expected_code
    upload_root = tmp_path / "uploads"
    assert not upload_root.exists() or not list(upload_root.rglob(filename))


def test_knowledge_library_upload_records_normalized_untrusted_format_metadata(tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "source_uploads_dir", str(tmp_path / "uploads"))
    payload = "网络设备配置基础\n交换机与路由器".encode("gb18030")
    db = _install_upload_endpoint_fakes(monkeypatch)
    response = asyncio.run(knowledge_library_api.upload_library_source(
        file=_upload_file("network-notes.txt", payload, "text/plain"),
        current=SimpleNamespace(learner=SimpleNamespace(id=9)),
        db=db,
    ))

    metadata = response["format"]
    assert metadata["registry_version"] == FORMAT_REGISTRY_VERSION
    assert metadata["format_id"] == "plain_text"
    assert metadata["encoding"] == "gb18030"
    assert metadata["mastery_inference"] is False
    assert metadata["executable"] is False
    assert "不可信" in metadata["trust_boundary"]
