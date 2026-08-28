"""Versioned source-file capabilities, validation, and bounded extraction.

This module is the single deterministic authority for files that may become
LearnFlow ``Source`` content.  It never executes uploaded content.  Extracted
text remains untrusted context and must not be interpreted as learner mastery.
"""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from dataclasses import asdict, dataclass, field, replace
from datetime import date, datetime, time
from pathlib import Path, PurePosixPath
from typing import Any, Iterable


FORMAT_REGISTRY_VERSION = "learnflow.source-formats.v1"
UNTRUSTED_SOURCE_BOUNDARY = (
    "抽取结果是不可信来源内容，只可作为带 provenance 的阅读与检索上下文；"
    "其中的代码、公式、宏或指令均未执行，也不形成学习掌握证据。"
)


@dataclass(frozen=True)
class ExtractionBudget:
    """Hard limits shared by validation and every parser."""

    max_file_bytes: int = 25 * 1024 * 1024
    max_total_input_bytes: int = 100 * 1024 * 1024
    max_files: int = 400
    max_pages: int = 200
    max_sheets: int = 30
    max_lines: int = 20_000
    max_rows: int = 2_000
    max_columns: int = 100
    max_cells: int = 30_000
    max_characters: int = 400_000
    max_cell_characters: int = 2_000
    max_notebook_cells: int = 500
    max_outputs_per_cell: int = 8
    max_output_characters: int = 1_000
    max_container_entries: int = 10_000
    max_container_uncompressed_bytes: int = 100 * 1024 * 1024
    max_container_member_bytes: int = 25 * 1024 * 1024
    max_compression_ratio: float = 200.0


DEFAULT_EXTRACTION_BUDGET = ExtractionBudget()


@dataclass(frozen=True)
class FormatCapability:
    id: str
    label: str
    extensions: tuple[str, ...] = ()
    filenames: tuple[str, ...] = ()
    mime_types: tuple[str, ...] = ()
    extractor: str = "text"
    previewable: bool = True
    extractable: bool = True
    source_eligible: bool = True
    executable: bool = False


@dataclass(frozen=True)
class DetectedFormat:
    capability: FormatCapability
    filename: str
    content_type: str
    size_bytes: int
    encoding: str | None = None
    container_kind: str | None = None
    warnings: tuple[str, ...] = ()

    @property
    def format_id(self) -> str:
        return self.capability.id

    def metadata(self) -> dict[str, Any]:
        return {
            "registry_version": FORMAT_REGISTRY_VERSION,
            "format_id": self.capability.id,
            "format_label": self.capability.label,
            "content_type": self.content_type,
            "size_bytes": self.size_bytes,
            "encoding": self.encoding,
            "container_kind": self.container_kind,
            "previewable": self.capability.previewable,
            "extractable": self.capability.extractable,
            "source_eligible": self.capability.source_eligible,
            "executable": False,
            "execution_performed": False,
            "warnings": list(self.warnings),
            "trust_boundary": UNTRUSTED_SOURCE_BOUNDARY,
            "mastery_inference": False,
        }


@dataclass(frozen=True)
class ExtractionResult:
    text: str
    detected: DetectedFormat
    truncated: bool = False
    warnings: tuple[str, ...] = ()
    counters: dict[str, int] = field(default_factory=dict)

    def metadata(self) -> dict[str, Any]:
        return {
            **self.detected.metadata(),
            "truncated": self.truncated,
            "warnings": [*self.detected.warnings, *self.warnings],
            "counters": dict(self.counters),
        }


class FileFormatError(ValueError):
    """An explainable, stable failure at the source-file boundary."""

    def __init__(self, code: str, message: str, *, status_code: int = 415):
        super().__init__(message)
        self.code = code
        self.status_code = status_code

    def detail(self) -> dict[str, str]:
        return {
            "code": self.code,
            "message": str(self),
            "format_registry_version": FORMAT_REGISTRY_VERSION,
        }


TEXT_MIME_TYPES = (
    "text/plain", "text/markdown", "text/x-markdown", "text/x-rst", "text/csv",
    "text/html", "text/css", "text/xml", "text/yaml", "text/x-yaml",
    "text/x-python", "text/x-c", "text/x-c++", "text/x-java-source",
    "text/x-shellscript", "text/x-sql", "text/javascript",
    "application/json", "application/ld+json", "application/jsonlines",
    "application/xml", "application/yaml", "application/x-yaml",
    "application/toml", "application/javascript", "application/sql",
)

SOURCE_CODE_EXTENSIONS = (
    ".py", ".pyi", ".js", ".mjs", ".cjs", ".jsx", ".ts", ".tsx",
    ".html", ".htm", ".css", ".scss", ".sass", ".less", ".vue", ".svelte",
    ".c", ".h", ".cc", ".cpp", ".cxx", ".hh", ".hpp", ".hxx",
    ".java", ".kt", ".kts", ".cs", ".go", ".rs", ".rb", ".php",
    ".jsp", ".jspx", ".asp", ".aspx", ".swift", ".dart", ".scala",
    ".groovy", ".pl", ".pm", ".vb", ".fs", ".fsx", ".sol", ".clj",
    ".cljs", ".ex", ".exs", ".erl", ".hrl", ".jl", ".sql", ".sh",
    ".bash", ".zsh", ".ps1", ".bat", ".cmd", ".lua", ".r", ".m",
    ".mm", ".asm", ".s",
)

CONFIG_EXTENSIONS = (
    ".json", ".jsonl", ".ndjson", ".xml", ".yaml", ".yml", ".toml",
    ".ini", ".cfg", ".conf", ".properties", ".gradle", ".graphql",
    ".gql", ".proto", ".tex", ".bib",
)

SPECIAL_TEXT_FILENAMES = (
    "dockerfile", "makefile", "gnumakefile", "cmakelists.txt", "requirements.txt",
    "constraints.txt", "pyproject.toml", "package.json", "package-lock.json",
    "pnpm-lock.yaml", "yarn.lock", "pom.xml", "build.gradle", "settings.gradle",
    "cargo.toml", "cargo.lock", "go.mod", "go.sum", "gemfile", "composer.json",
    "pipfile", "pipfile.lock", "tox.ini", "pytest.ini", "mypy.ini",
    ".gitignore", ".dockerignore", ".editorconfig", ".npmrc", ".prettierrc",
    ".eslintrc",
)

FORMAT_CAPABILITIES: tuple[FormatCapability, ...] = (
    FormatCapability(
        "plain_text", "纯文本", (".txt", ".log"), mime_types=TEXT_MIME_TYPES,
    ),
    FormatCapability(
        "markdown", "Markdown", (".md", ".markdown", ".mdown", ".mkd"),
        mime_types=TEXT_MIME_TYPES,
    ),
    FormatCapability(
        "rst", "reStructuredText", (".rst",), mime_types=TEXT_MIME_TYPES,
    ),
    FormatCapability(
        "csv", "CSV 表格", (".csv",), mime_types=("text/csv", "text/plain"),
        extractor="csv",
    ),
    FormatCapability(
        "pdf", "PDF 文档", (".pdf",), mime_types=("application/pdf",),
        extractor="pdf",
    ),
    FormatCapability(
        "docx", "Word 文档", (".docx",),
        mime_types=("application/vnd.openxmlformats-officedocument.wordprocessingml.document",),
        extractor="docx",
    ),
    FormatCapability(
        "pptx", "PowerPoint 演示文稿", (".pptx",),
        mime_types=("application/vnd.openxmlformats-officedocument.presentationml.presentation",),
        extractor="pptx",
    ),
    FormatCapability(
        "xlsx", "Excel 工作簿", (".xlsx",),
        mime_types=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",),
        extractor="xlsx",
    ),
    FormatCapability(
        "ipynb", "Jupyter Notebook", (".ipynb",),
        mime_types=("application/x-ipynb+json", "application/json"),
        extractor="ipynb",
    ),
    FormatCapability(
        "source_code", "源代码", SOURCE_CODE_EXTENSIONS,
        filenames=("dockerfile", "makefile", "gnumakefile", "cmakelists.txt"),
        mime_types=TEXT_MIME_TYPES,
    ),
    FormatCapability(
        "configuration", "配置与结构化文本", CONFIG_EXTENSIONS,
        filenames=tuple(name for name in SPECIAL_TEXT_FILENAMES if name not in {
            "dockerfile", "makefile", "gnumakefile", "cmakelists.txt",
        }),
        mime_types=TEXT_MIME_TYPES,
    ),
)

_CAPABILITY_BY_EXTENSION = {
    extension: capability
    for capability in FORMAT_CAPABILITIES
    for extension in capability.extensions
}
_CAPABILITY_BY_FILENAME = {
    filename: capability
    for capability in FORMAT_CAPABILITIES
    for filename in capability.filenames
}

MACRO_OFFICE_EXTENSIONS = {
    ".docm", ".dotm", ".xlsm", ".xltm", ".xlam",
    ".pptm", ".potm", ".ppam", ".ppsm", ".sldm",
}
ARCHIVE_EXTENSIONS = {
    ".zip", ".rar", ".7z", ".tar", ".tgz", ".gz", ".bz2", ".xz",
    ".jar", ".war", ".ear", ".apk", ".aab",
}
BINARY_EXECUTABLE_EXTENSIONS = {
    ".exe", ".dll", ".com", ".msi", ".so", ".dylib", ".app", ".deb",
    ".rpm", ".bin", ".class", ".wasm",
}
SECRET_EXTENSIONS = {".pem", ".key", ".p12", ".pfx", ".keystore", ".jks"}
SECRET_FILENAMES = {".env", ".env.local", ".env.production", "credentials.json"}

GENERIC_MIME_TYPES = {"", "application/octet-stream", "binary/octet-stream"}


def registry_manifest() -> dict[str, Any]:
    """Return a machine-readable, immutable view used by tests and adapters."""

    return {
        "version": FORMAT_REGISTRY_VERSION,
        "trust_boundary": UNTRUSTED_SOURCE_BOUNDARY,
        "formats": [asdict(capability) for capability in FORMAT_CAPABILITIES],
        "budgets": asdict(DEFAULT_EXTRACTION_BUDGET),
        "explicitly_unsupported": {
            "macro_office": sorted(MACRO_OFFICE_EXTENSIONS),
            "archives": sorted(ARCHIVE_EXTENSIONS),
            "binary_executables": sorted(BINARY_EXECUTABLE_EXTENSIONS),
            "audio_video_transcription": False,
            "safe_execution": False,
        },
    }


def source_extensions() -> frozenset[str]:
    return frozenset(_CAPABILITY_BY_EXTENSION)


def capability_for_filename(filename: str) -> FormatCapability | None:
    basename = Path(str(filename).replace("\\", "/")).name.casefold()
    if basename in _CAPABILITY_BY_FILENAME:
        return _CAPABILITY_BY_FILENAME[basename]
    return _CAPABILITY_BY_EXTENSION.get(Path(basename).suffix.casefold())


def format_id_for_filename(filename: str) -> str | None:
    capability = capability_for_filename(filename)
    return capability.id if capability else None


def is_source_filename(filename: str) -> bool:
    return capability_for_filename(filename) is not None


def _normalized_mime(content_type: str | None) -> str:
    return str(content_type or "").split(";", 1)[0].strip().casefold()


def _known_binary_magic(data: bytes) -> str | None:
    if data.startswith(b"%PDF-"):
        return "pdf"
    if data.startswith((b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08")):
        return "zip"
    if data.startswith(b"\x7fELF"):
        return "elf_executable"
    if data.startswith(b"MZ"):
        return "pe_executable"
    if data[:4] in {
        b"\xfe\xed\xfa\xce", b"\xce\xfa\xed\xfe", b"\xfe\xed\xfa\xcf",
        b"\xcf\xfa\xed\xfe", b"\xca\xfe\xba\xbe", b"\xbe\xba\xfe\xca",
    }:
        return "mach_executable"
    if data.startswith(b"Rar!"):
        return "rar_archive"
    if data.startswith(b"7z\xbc\xaf\x27\x1c"):
        return "7z_archive"
    if data.startswith(b"\x1f\x8b"):
        return "gzip_archive"
    if data.startswith(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"):
        return "ole_compound"
    if len(data) > 262 and data[257:262] == b"ustar":
        return "tar_archive"
    return None


def _decode_text(data: bytes) -> tuple[str, str]:
    if b"\x00" in data:
        raise FileFormatError("binary_content", "文件包含 NUL 字节，不能作为文本来源")
    if data.startswith(b"\xef\xbb\xbf"):
        attempts = (("utf-8-sig", "utf-8-bom"),)
    else:
        attempts = (("utf-8", "utf-8"), ("gb18030", "gb18030"))
    decoded = None
    label = ""
    for codec, candidate_label in attempts:
        try:
            decoded = data.decode(codec, errors="strict")
            label = candidate_label
            break
        except UnicodeDecodeError:
            continue
    if decoded is None:
        raise FileFormatError(
            "unsupported_text_encoding",
            "文本编码无法识别；当前只支持 UTF-8、UTF-8 BOM 和 GB18030",
        )
    controls = sum(
        1 for character in decoded
        if ord(character) < 32 and character not in "\n\r\t\f"
    )
    if decoded and controls / max(len(decoded), 1) > 0.01:
        raise FileFormatError("binary_content", "文件含有过多控制字符，疑似二进制内容")
    return decoded, label


def _reject_by_name(filename: str) -> None:
    raw_basename = Path(str(filename).replace("\\", "/")).name
    if any(ord(character) < 32 for character in raw_basename):
        raise FileFormatError("unsafe_filename", "文件名包含控制字符", status_code=400)
    if len(raw_basename.encode("utf-8")) > 255:
        raise FileFormatError("unsafe_filename", "文件名超过 255 字节", status_code=400)
    basename = raw_basename.casefold()
    suffix = Path(basename).suffix.casefold()
    if suffix in MACRO_OFFICE_EXTENSIONS:
        raise FileFormatError("macro_office_rejected", f"拒绝含宏 Office 文件：{suffix}")
    if suffix in ARCHIVE_EXTENSIONS or basename.endswith((".tar.gz", ".tar.bz2", ".tar.xz")):
        raise FileFormatError("archive_rejected", f"压缩包不能直接作为学习来源：{basename}")
    if suffix in BINARY_EXECUTABLE_EXTENSIONS:
        raise FileFormatError("binary_executable_rejected", f"未知可执行文件不能作为学习来源：{suffix}")
    if suffix in SECRET_EXTENSIONS or basename in SECRET_FILENAMES or basename.startswith(".env."):
        raise FileFormatError("secret_file_rejected", "密钥或环境凭据文件不能导入领域知识库")


def _validate_mime(capability: FormatCapability, content_type: str) -> None:
    mime = _normalized_mime(content_type)
    if mime in GENERIC_MIME_TYPES:
        return
    if mime == "application/zip" and capability.extractor in {"docx", "pptx", "xlsx"}:
        return
    allowed = {item.casefold() for item in capability.mime_types}
    if mime not in allowed:
        raise FileFormatError(
            "mime_extension_mismatch",
            f"文件扩展名对应 {capability.label}，但上传 MIME 为 {mime or '空'}",
        )


def validate_declared_format(
    filename: str,
    content_type: str | None = None,
) -> FormatCapability:
    """Validate the user-controlled name/MIME pair before storing its bytes."""

    _reject_by_name(filename)
    capability = capability_for_filename(filename)
    if capability is None:
        raise FileFormatError(
            "unsupported_format",
            f"不支持的来源文件格式：{Path(filename).suffix or filename}",
        )
    _validate_mime(capability, content_type or "")
    return capability


def _safe_zip_member(name: str) -> bool:
    if not name or "\\" in name or name.startswith("/"):
        return False
    path = PurePosixPath(name)
    return not path.is_absolute() and ".." not in path.parts


def _validate_ooxml_container(
    data: bytes,
    expected_kind: str,
    budget: ExtractionBudget,
) -> tuple[str, tuple[str, ...]]:
    try:
        archive = zipfile.ZipFile(io.BytesIO(data))
    except (zipfile.BadZipFile, OSError) as exc:
        raise FileFormatError("corrupt_ooxml", "Office 文件不是有效的 OOXML 容器") from exc
    with archive:
        entries = archive.infolist()
        if len(entries) > budget.max_container_entries:
            raise FileFormatError("container_budget_exceeded", "Office 容器条目数量超过安全预算")
        total_uncompressed = 0
        names = set()
        for entry in entries:
            normalized = entry.filename.casefold()
            if not _safe_zip_member(entry.filename):
                raise FileFormatError("unsafe_container_path", "Office 容器包含不安全路径")
            mode = (entry.external_attr >> 16) & 0o170000
            if mode == 0o120000:
                raise FileFormatError("unsafe_container_link", "Office 容器包含符号链接")
            if normalized in names:
                raise FileFormatError("duplicate_container_member", "Office 容器包含重复成员")
            if entry.file_size > budget.max_container_member_bytes:
                raise FileFormatError("container_budget_exceeded", "Office 容器单个成员超过安全预算")
            total_uncompressed += entry.file_size
            if total_uncompressed > budget.max_container_uncompressed_bytes:
                raise FileFormatError("container_budget_exceeded", "Office 容器展开大小超过安全预算")
            ratio = entry.file_size / max(entry.compress_size, 1)
            if entry.file_size > 1_000_000 and ratio > budget.max_compression_ratio:
                raise FileFormatError("suspicious_compression_ratio", "Office 容器压缩比异常")
            names.add(normalized)
        if any(
            "vbaproject.bin" in name or "/macrosheets/" in f"/{name}"
            for name in names
        ):
            raise FileFormatError("macro_office_rejected", "Office 容器包含 VBA/XLM 宏，已拒绝")
        try:
            content_types = archive.read("[Content_Types].xml").lower()
        except KeyError as exc:
            raise FileFormatError("corrupt_ooxml", "Office 容器缺少 [Content_Types].xml") from exc
        if any(marker in content_types for marker in (b"macroenabled", b"vbaproject", b"macrosheet")):
            raise FileFormatError("macro_office_rejected", "Office 容器声明了宏内容，已拒绝")

    kind_by_required_member = {
        "docx": "word/document.xml",
        "pptx": "ppt/presentation.xml",
        "xlsx": "xl/workbook.xml",
    }
    actual_kinds = [
        kind for kind, required in kind_by_required_member.items()
        if required in names
    ]
    if actual_kinds != [expected_kind]:
        actual = actual_kinds[0] if actual_kinds else "unknown"
        raise FileFormatError(
            "ooxml_extension_mismatch",
            f"文件扩展名声明 {expected_kind}，但容器内容识别为 {actual}",
        )
    return expected_kind, ()


def inspect_bytes(
    data: bytes,
    filename: str,
    content_type: str | None = None,
    *,
    budget: ExtractionBudget = DEFAULT_EXTRACTION_BUDGET,
) -> DetectedFormat:
    """Validate extension, bounded MIME/magic, encoding, and OOXML structure."""

    try:
        capability = validate_declared_format(filename, content_type)
    except FileFormatError as declared_error:
        magic = _known_binary_magic(data)
        if declared_error.code == "unsupported_format":
            if magic and "executable" in magic:
                raise FileFormatError(
                    "binary_executable_rejected",
                    "检测到未知二进制可执行文件，已拒绝",
                )
            if magic and ("archive" in magic or magic == "zip"):
                raise FileFormatError("archive_rejected", "检测到压缩包，不能直接作为学习来源")
            if magic == "ole_compound":
                raise FileFormatError("legacy_office_rejected", "旧版 Office/OLE 文件不在安全来源范围内")
        raise
    if not data:
        raise FileFormatError("empty_file", "文件为空，无法建立学习来源", status_code=400)
    if len(data) > budget.max_file_bytes:
        raise FileFormatError("file_budget_exceeded", "文件大小超过格式解析预算", status_code=413)
    magic = _known_binary_magic(data)
    encoding = None
    container_kind = None
    if capability.extractor == "pdf":
        if magic and "executable" in magic:
            raise FileFormatError("binary_executable_rejected", "PDF 扩展名伪装了二进制可执行文件")
        if magic and ("archive" in magic or magic == "zip"):
            raise FileFormatError("archive_rejected", "PDF 扩展名伪装了压缩包")
        if magic != "pdf":
            raise FileFormatError("magic_extension_mismatch", "扩展名是 PDF，但文件头不是 %PDF-")
    elif capability.extractor in {"docx", "pptx", "xlsx"}:
        if magic != "zip":
            raise FileFormatError("magic_extension_mismatch", "OOXML 文件缺少有效 ZIP 文件头")
        container_kind, _ = _validate_ooxml_container(data, capability.extractor, budget)
    else:
        if magic is not None:
            if "executable" in magic:
                raise FileFormatError("binary_executable_rejected", "文本扩展名伪装了二进制可执行文件")
            if magic in {"zip", "rar_archive", "7z_archive", "gzip_archive", "tar_archive"}:
                raise FileFormatError("archive_rejected", "文本扩展名伪装了压缩包")
            if magic == "ole_compound":
                raise FileFormatError("legacy_office_rejected", "文本扩展名伪装了旧版 Office/OLE 文件")
            raise FileFormatError("magic_extension_mismatch", "文件扩展名与检测到的二进制格式不一致")
        decoded, encoding = _decode_text(data)
        if capability.extractor == "ipynb":
            try:
                notebook = json.loads(decoded)
            except json.JSONDecodeError as exc:
                raise FileFormatError("corrupt_ipynb", f"Notebook JSON 损坏：第 {exc.lineno} 行") from exc
            if not isinstance(notebook, dict) or not isinstance(notebook.get("cells", []), list):
                raise FileFormatError("corrupt_ipynb", "Notebook 缺少合法 cells 数组")

    return DetectedFormat(
        capability=capability,
        filename=Path(str(filename).replace("\\", "/")).name,
        content_type=_normalized_mime(content_type),
        size_bytes=len(data),
        encoding=encoding,
        container_kind=container_kind,
    )


def inspect_path(
    path: str | Path,
    *,
    filename: str | None = None,
    content_type: str | None = None,
    budget: ExtractionBudget = DEFAULT_EXTRACTION_BUDGET,
) -> DetectedFormat:
    target = Path(path)
    try:
        size = target.stat().st_size
    except OSError as exc:
        raise FileFormatError("file_unreadable", "无法读取上传文件", status_code=400) from exc
    if size > budget.max_file_bytes:
        raise FileFormatError("file_budget_exceeded", "文件大小超过格式解析预算", status_code=413)
    try:
        data = target.read_bytes()
    except OSError as exc:
        raise FileFormatError("file_unreadable", "无法读取上传文件", status_code=400) from exc
    return inspect_bytes(data, filename or target.name, content_type, budget=budget)


class _BoundedText:
    def __init__(self, budget: ExtractionBudget):
        self.budget = budget
        self.parts: list[str] = []
        self.characters = 0
        self.lines = 0
        self.truncated = False

    def add(self, value: Any) -> None:
        if self.truncated:
            return
        text = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
        if not text:
            return
        separator = "\n\n" if self.parts else ""
        remaining_lines = self.budget.max_lines - self.lines
        separator_line_cost = 1 if separator else 0
        allowed_text_lines = remaining_lines - separator_line_cost
        if allowed_text_lines <= 0:
            self.truncated = True
            return
        lines = text.splitlines()
        if len(lines) > allowed_text_lines:
            lines = lines[:allowed_text_lines]
            text = "\n".join(lines)
            self.truncated = True
        remaining_chars = self.budget.max_characters - self.characters
        if remaining_chars <= len(separator):
            self.truncated = True
            return
        if len(text) > remaining_chars - len(separator):
            text = text[:remaining_chars - len(separator)]
            self.truncated = True
        if not text:
            return
        self.parts.append(f"{separator}{text}")
        self.characters += len(separator) + len(text)
        self.lines += separator_line_cost + text.count("\n") + 1

    def text(self) -> str:
        return "".join(self.parts).strip()


def _bounded_cell(value: Any, budget: ExtractionBudget) -> tuple[str, bool]:
    if value is None:
        return "", False
    if isinstance(value, (datetime, date, time)):
        text = value.isoformat()
    else:
        text = str(value)
    text = text.replace("\r", " ").replace("\n", " ")
    if len(text) > budget.max_cell_characters:
        return text[:budget.max_cell_characters], True
    return text, False


def _extract_text(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    decoded, encoding = _decode_text(data)
    builder = _BoundedText(budget)
    builder.add(decoded)
    if not builder.text():
        raise FileFormatError("no_extractable_text", "文件没有可抽取的文本内容", status_code=400)
    actual = replace(detected, encoding=encoding)
    return ExtractionResult(
        builder.text(), actual, builder.truncated,
        counters={"characters": builder.characters, "lines": builder.lines},
    )


def _extract_csv(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    decoded, encoding = _decode_text(data)
    sample = decoded[:32_768]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",\t;|")
    except csv.Error:
        dialect = csv.excel
    builder = _BoundedText(budget)
    builder.add(f"# {detected.filename}")
    rows = cells = 0
    truncated = False
    try:
        reader = csv.reader(io.StringIO(decoded), dialect, strict=True)
        for row_number, row in enumerate(reader, start=1):
            if row_number > budget.max_rows or cells >= budget.max_cells:
                truncated = True
                break
            if len(row) > budget.max_columns:
                row = row[:budget.max_columns]
                truncated = True
            normalized = []
            for value in row:
                if cells >= budget.max_cells:
                    truncated = True
                    break
                cell, cell_truncated = _bounded_cell(value, budget)
                normalized.append(cell)
                cells += 1
                truncated = truncated or cell_truncated
            builder.add(f"[行 {row_number}] {json.dumps(normalized, ensure_ascii=False)}")
            rows += 1
            if builder.truncated:
                truncated = True
                break
    except csv.Error as exc:
        raise FileFormatError("corrupt_csv", f"CSV 解析失败：{exc}", status_code=400) from exc
    if rows == 0:
        raise FileFormatError("no_extractable_text", "CSV 没有可抽取的行", status_code=400)
    actual = replace(detected, encoding=encoding)
    return ExtractionResult(
        builder.text(), actual, truncated or builder.truncated,
        counters={"rows": rows, "cells": cells, "characters": builder.characters},
    )


def _extract_pdf(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(data), strict=False)
    except Exception as exc:
        raise FileFormatError("corrupt_pdf", f"PDF 解析失败：{type(exc).__name__}", status_code=400) from exc
    if reader.is_encrypted:
        raise FileFormatError("encrypted_document", "加密 PDF 不能作为学习来源")
    builder = _BoundedText(budget)
    extracted_pages = 0
    text_pages = 0
    warnings = []
    for page_number, page in enumerate(reader.pages, start=1):
        if page_number > budget.max_pages:
            break
        try:
            content = page.extract_text() or ""
        except Exception as exc:
            warnings.append(f"第 {page_number} 页抽取失败：{type(exc).__name__}")
            content = ""
        builder.add(f"## 第 {page_number} 页")
        if content.strip():
            builder.add(content)
            text_pages += 1
        else:
            builder.add("[本页没有可抽取文本]")
        extracted_pages += 1
        if builder.truncated:
            break
    if text_pages == 0:
        raise FileFormatError(
            "no_extractable_text",
            "PDF 没有可抽取文本；扫描件 OCR 尚未支持",
            status_code=400,
        )
    truncated = len(reader.pages) > extracted_pages or builder.truncated
    return ExtractionResult(
        builder.text(), detected, truncated, tuple(warnings),
        {"pages": extracted_pages, "text_pages": text_pages, "characters": builder.characters},
    )


def _iter_docx_blocks(document: Any) -> Iterable[Any]:
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P

    for child in document.element.body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _heading_level(paragraph: Any) -> int | None:
    style = getattr(paragraph, "style", None)
    value = " ".join(filter(None, [
        str(getattr(style, "name", "") or ""),
        str(getattr(style, "style_id", "") or ""),
    ]))
    match = re.search(r"(?:heading|标题)\s*([1-6])", value, re.IGNORECASE)
    return int(match.group(1)) if match else None


def _extract_docx(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    try:
        from docx import Document
        from docx.table import Table

        document = Document(io.BytesIO(data))
    except Exception as exc:
        raise FileFormatError("corrupt_docx", f"DOCX 解析失败：{type(exc).__name__}", status_code=400) from exc
    builder = _BoundedText(budget)
    paragraphs = tables = rows = cells = 0
    truncated = False
    for block in _iter_docx_blocks(document):
        if isinstance(block, Table):
            tables += 1
            builder.add(f"## 表格 {tables}")
            for row_number, row in enumerate(block.rows, start=1):
                if rows >= budget.max_rows or cells >= budget.max_cells:
                    truncated = True
                    break
                values = []
                for cell in row.cells[:budget.max_columns]:
                    if cells >= budget.max_cells:
                        truncated = True
                        break
                    value, cell_truncated = _bounded_cell(cell.text, budget)
                    values.append(value)
                    cells += 1
                    truncated = truncated or cell_truncated
                if len(row.cells) > budget.max_columns:
                    truncated = True
                builder.add(f"[行 {row_number}] {json.dumps(values, ensure_ascii=False)}")
                rows += 1
                if builder.truncated:
                    truncated = True
                    break
        else:
            text = str(block.text or "").strip()
            if not text:
                continue
            level = _heading_level(block)
            builder.add(f"{'#' * level} {text}" if level else text)
            paragraphs += 1
        if builder.truncated or truncated and cells >= budget.max_cells:
            break
    if not builder.text():
        raise FileFormatError("no_extractable_text", "DOCX 没有可抽取的标题、段落或表格", status_code=400)
    return ExtractionResult(
        builder.text(), detected, truncated or builder.truncated,
        counters={
            "paragraphs": paragraphs, "tables": tables, "rows": rows,
            "cells": cells, "characters": builder.characters,
        },
    )


def _extract_pptx(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise FileFormatError("corrupt_pptx", f"PPTX 解析失败：{type(exc).__name__}", status_code=400) from exc
    builder = _BoundedText(budget)
    slides = notes = rows = cells = 0
    truncated = False
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if slide_number > budget.max_pages:
            truncated = True
            break
        title_shape = slide.shapes.title
        title = str(title_shape.text or "").strip() if title_shape else ""
        builder.add(f"## 幻灯片 {slide_number}{f'：{title}' if title else ''}")
        for shape in slide.shapes:
            if title_shape is not None and shape.shape_id == title_shape.shape_id:
                continue
            if getattr(shape, "has_text_frame", False):
                text = str(shape.text or "").strip()
                if text:
                    builder.add(text)
            if getattr(shape, "has_table", False):
                builder.add(f"### 幻灯片 {slide_number} 表格")
                for row_number, row in enumerate(shape.table.rows, start=1):
                    if rows >= budget.max_rows or cells >= budget.max_cells:
                        truncated = True
                        break
                    values = []
                    for column_number, cell in enumerate(row.cells, start=1):
                        if column_number > budget.max_columns:
                            truncated = True
                            break
                        if cells >= budget.max_cells:
                            truncated = True
                            break
                        value, cell_truncated = _bounded_cell(cell.text, budget)
                        values.append(value)
                        cells += 1
                        truncated = truncated or cell_truncated
                    builder.add(f"[行 {row_number}] {json.dumps(values, ensure_ascii=False)}")
                    rows += 1
        try:
            notes_text = str(slide.notes_slide.notes_text_frame.text or "").strip()
        except (AttributeError, KeyError, ValueError):
            notes_text = ""
        if notes_text:
            builder.add(f"### 备注\n{notes_text}")
            notes += 1
        slides += 1
        if builder.truncated:
            truncated = True
            break
    if not builder.text():
        raise FileFormatError(
            "no_extractable_text",
            "PPTX 没有可抽取的标题、正文、表格或备注",
            status_code=400,
        )
    return ExtractionResult(
        builder.text(), detected, truncated or builder.truncated,
        counters={
            "slides": slides, "notes": notes, "rows": rows, "cells": cells,
            "characters": builder.characters,
        },
    )


def _extract_xlsx(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(
            io.BytesIO(data), read_only=True, data_only=False, keep_links=False,
        )
    except Exception as exc:
        raise FileFormatError("corrupt_xlsx", f"XLSX 解析失败：{type(exc).__name__}", status_code=400) from exc
    builder = _BoundedText(budget)
    sheets = rows = cells = formulas = 0
    truncated = False
    try:
        for sheet_number, worksheet in enumerate(workbook.worksheets, start=1):
            if sheet_number > budget.max_sheets:
                truncated = True
                break
            builder.add(f"## 工作表：{worksheet.title}")
            max_row = min(int(worksheet.max_row or 0), budget.max_rows)
            max_column = min(int(worksheet.max_column or 0), budget.max_columns)
            if int(worksheet.max_row or 0) > budget.max_rows or int(worksheet.max_column or 0) > budget.max_columns:
                truncated = True
            for row_number, row in enumerate(
                worksheet.iter_rows(
                    min_row=1, max_row=max_row, max_col=max_column, values_only=True,
                ),
                start=1,
            ):
                if cells >= budget.max_cells:
                    truncated = True
                    break
                values = []
                nonempty = False
                for value in row:
                    if cells >= budget.max_cells:
                        truncated = True
                        break
                    normalized, cell_truncated = _bounded_cell(value, budget)
                    values.append(normalized)
                    nonempty = nonempty or bool(normalized)
                    formulas += int(normalized.startswith("="))
                    cells += 1
                    truncated = truncated or cell_truncated
                if nonempty:
                    builder.add(f"[行 {row_number}] {json.dumps(values, ensure_ascii=False)}")
                    rows += 1
                if builder.truncated:
                    truncated = True
                    break
            sheets += 1
            if builder.truncated:
                break
    finally:
        workbook.close()
    if rows == 0:
        raise FileFormatError("no_extractable_text", "XLSX 没有可抽取的单元格值或公式", status_code=400)
    return ExtractionResult(
        builder.text(), detected, truncated or builder.truncated,
        counters={
            "sheets": sheets, "rows": rows, "cells": cells, "formulas": formulas,
            "characters": builder.characters,
        },
    )


_SENSITIVE_PATTERNS = (
    re.compile(r"(?i)\b(api[_-]?key|token|password|secret)\b\s*[:=]\s*([^\s,;]+)"),
    re.compile(r"\bsk-[A-Za-z0-9_-]{12,}\b"),
    re.compile(r"\bAKIA[0-9A-Z]{16}\b"),
    re.compile(r"-----BEGIN [A-Z ]+ PRIVATE KEY-----[\s\S]*?-----END [A-Z ]+ PRIVATE KEY-----"),
)


def _redact_notebook_text(value: Any, limit: int) -> tuple[str, bool]:
    text = "".join(value) if isinstance(value, list) else str(value or "")
    for pattern in _SENSITIVE_PATTERNS:
        text = pattern.sub(lambda match: f"{match.group(1)}=<redacted>" if match.lastindex else "<redacted>", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    if len(text) > limit:
        return text[:limit], True
    return text, False


def _extract_ipynb(data: bytes, detected: DetectedFormat, budget: ExtractionBudget) -> ExtractionResult:
    decoded, encoding = _decode_text(data)
    try:
        notebook = json.loads(decoded)
    except json.JSONDecodeError as exc:
        raise FileFormatError("corrupt_ipynb", f"Notebook JSON 损坏：第 {exc.lineno} 行", status_code=400) from exc
    cells_data = notebook.get("cells", [])
    if not isinstance(cells_data, list):
        raise FileFormatError("corrupt_ipynb", "Notebook cells 必须是数组", status_code=400)
    builder = _BoundedText(budget)
    metadata = notebook.get("metadata") if isinstance(notebook.get("metadata"), dict) else {}
    kernelspec = metadata.get("kernelspec") if isinstance(metadata.get("kernelspec"), dict) else {}
    language_info = metadata.get("language_info") if isinstance(metadata.get("language_info"), dict) else {}
    language = str(language_info.get("name") or kernelspec.get("language") or "").casefold()
    language = language if re.fullmatch(r"[a-z0-9_+.-]{1,32}", language) else ""
    cells = outputs = 0
    truncated = False
    for cell_number, cell in enumerate(cells_data, start=1):
        if cell_number > budget.max_notebook_cells:
            truncated = True
            break
        if not isinstance(cell, dict):
            continue
        cell_type = str(cell.get("cell_type") or "unknown")
        source_value = (
            "".join(cell.get("source") or [])
            if isinstance(cell.get("source"), list)
            else str(cell.get("source") or "")
        )
        source, source_clipped = _redact_notebook_text(source_value, budget.max_cell_characters)
        truncated = truncated or source_clipped
        if cell_type == "markdown":
            builder.add(f"## Markdown 单元 {cell_number}\n{source}")
        elif cell_type == "code":
            builder.add(f"## Code 单元 {cell_number}\n```{language}\n{source}\n```")
            cell_outputs = cell.get("outputs") if isinstance(cell.get("outputs"), list) else []
            for output_number, output in enumerate(cell_outputs, start=1):
                if output_number > budget.max_outputs_per_cell:
                    truncated = True
                    break
                if not isinstance(output, dict):
                    continue
                output_type = str(output.get("output_type") or "unknown")
                summary = ""
                if output_type == "stream":
                    summary, clipped = _redact_notebook_text(output.get("text", ""), budget.max_output_characters)
                elif output_type == "error":
                    summary, clipped = _redact_notebook_text(
                        f"{output.get('ename', 'Error')}: {output.get('evalue', '')}",
                        budget.max_output_characters,
                    )
                elif output_type in {"execute_result", "display_data"}:
                    payload = output.get("data") if isinstance(output.get("data"), dict) else {}
                    text_value = payload.get("text/plain", "")
                    summary, clipped = _redact_notebook_text(text_value, budget.max_output_characters)
                    omitted = sorted(key for key in payload if key != "text/plain")
                    if omitted:
                        summary = f"{summary}\n[省略富媒体输出：{', '.join(omitted[:12])}]".strip()
                else:
                    summary, clipped = f"[{output_type} 输出未展开]", False
                builder.add(f"### 输出 {output_number}（{output_type}）\n{summary}")
                outputs += 1
                truncated = truncated or clipped
        else:
            builder.add(f"## {cell_type} 单元 {cell_number}\n[未展开非 Markdown/Code 单元]")
        cells += 1
        if builder.truncated:
            truncated = True
            break
    if cells == 0:
        raise FileFormatError("no_extractable_text", "Notebook 没有可抽取单元", status_code=400)
    actual = replace(detected, encoding=encoding)
    return ExtractionResult(
        builder.text(), actual, truncated or builder.truncated,
        counters={"cells": cells, "outputs": outputs, "characters": builder.characters},
    )


_EXTRACTORS = {
    "text": _extract_text,
    "csv": _extract_csv,
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "xlsx": _extract_xlsx,
    "ipynb": _extract_ipynb,
}


def extract_bytes(
    data: bytes,
    filename: str,
    content_type: str | None = None,
    *,
    budget: ExtractionBudget = DEFAULT_EXTRACTION_BUDGET,
) -> ExtractionResult:
    detected = inspect_bytes(data, filename, content_type, budget=budget)
    extractor = _EXTRACTORS[detected.capability.extractor]
    try:
        return extractor(data, detected, budget)
    except FileFormatError:
        raise
    except Exception as exc:
        raise FileFormatError(
            f"{detected.format_id}_extraction_failed",
            f"{detected.capability.label}解析失败：{type(exc).__name__}",
            status_code=400,
        ) from exc


def extract_path(
    path: str | Path,
    *,
    filename: str | None = None,
    content_type: str | None = None,
    budget: ExtractionBudget = DEFAULT_EXTRACTION_BUDGET,
) -> ExtractionResult:
    target = Path(path)
    try:
        size = target.stat().st_size
        if size > budget.max_file_bytes:
            raise FileFormatError("file_budget_exceeded", "文件大小超过格式解析预算", status_code=413)
        data = target.read_bytes()
    except OSError as exc:
        raise FileFormatError("file_unreadable", "无法读取上传文件", status_code=400) from exc
    return extract_bytes(
        data,
        filename or target.name,
        content_type,
        budget=budget,
    )
